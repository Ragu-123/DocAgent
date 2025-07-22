import gradio as gr
import os
import json
import uuid
import asyncio
from datetime import datetime
from typing import List, Dict, Any, Optional, Generator
import logging

# Import required libraries
from huggingface_hub import InferenceClient
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document

# Import document parsers
import PyPDF2
from pptx import Presentation
import pandas as pd
from docx import Document as DocxDocument
import io

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Get HuggingFace token from environment
HF_TOKEN = os.getenv("hf_token")
if not HF_TOKEN:
    raise ValueError("HuggingFace token not found in environment variables")

# Initialize HuggingFace Inference Client
client = InferenceClient(model="meta-llama/Llama-3.1-8B-Instruct", token=HF_TOKEN)

# Initialize embeddings
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

class MCPMessage:
    """Model Context Protocol Message Structure"""
    def __init__(self, sender: str, receiver: str, msg_type: str,
                 trace_id: str = None, payload: Dict = None):
        self.sender = sender
        self.receiver = receiver
        self.type = msg_type
        self.trace_id = trace_id or str(uuid.uuid4())
        self.payload = payload or {}
        self.timestamp = datetime.now().isoformat()
    
    def to_dict(self):
        return {
            "sender": self.sender,
            "receiver": self.receiver,
            "type": self.type,
            "trace_id": self.trace_id,
            "payload": self.payload,
            "timestamp": self.timestamp
        }

class MessageBus:
    """In-memory message bus for MCP communication"""
    def __init__(self):
        self.messages = []
        self.subscribers = {}
    
    def publish(self, message: MCPMessage):
        """Publish message to the bus"""
        self.messages.append(message)
        logger.info(f"Message published: {message.sender} -> {message.receiver} [{message.type}]")
        
        # Notify subscribers
        if message.receiver in self.subscribers:
            for callback in self.subscribers[message.receiver]:
                callback(message)
    
    def subscribe(self, agent_name: str, callback):
        """Subscribe agent to receive messages"""
        if agent_name not in self.subscribers:
            self.subscribers[agent_name] = []
        self.subscribers[agent_name].append(callback)

# Global message bus
message_bus = MessageBus()

class IngestionAgent:
    """Agent responsible for document parsing and preprocessing"""
    
    def __init__(self, message_bus: MessageBus):
        self.name = "IngestionAgent"
        self.message_bus = message_bus
        self.message_bus.subscribe(self.name, self.handle_message)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
    
    def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        if message.type == "INGESTION_REQUEST":
            self.process_documents(message)
    
    def parse_pdf(self, file_path: str) -> str:
        """Parse PDF document"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                return text
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            return ""
    
    def parse_pptx(self, file_path: str) -> str:
        """Parse PPTX document"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            return ""
    
    def parse_csv(self, file_path: str) -> str:
        """Parse CSV document"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Error parsing CSV: {e}")
            return ""
    
    def parse_docx(self, file_path: str) -> str:
        """Parse DOCX document"""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            return ""
    
    def parse_txt(self, file_path: str) -> str:
        """Parse TXT/Markdown document"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error parsing TXT: {e}")
            return ""
    
    def process_documents(self, message: MCPMessage):
        """Process uploaded documents"""
        files = message.payload.get("files", [])
        processed_docs = []
        
        for file_path in files:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # Parse document based on file type
            if file_ext == '.pdf':
                text = self.parse_pdf(file_path)
            elif file_ext == '.pptx':
                text = self.parse_pptx(file_path)
            elif file_ext == '.csv':
                text = self.parse_csv(file_path)
            elif file_ext == '.docx':
                text = self.parse_docx(file_path)
            elif file_ext in ['.txt', '.md']:
                text = self.parse_txt(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_ext}")
                continue
            
            if text:
                # Split text into chunks
                chunks = self.text_splitter.split_text(text)
                docs = [Document(page_content=chunk, metadata={"source": file_path}) 
                        for chunk in chunks]
                processed_docs.extend(docs)
        
        # Send processed documents to RetrievalAgent
        response = MCPMessage(
            sender=self.name,
            receiver="RetrievalAgent",
            msg_type="INGESTION_COMPLETE",
            trace_id=message.trace_id,
            payload={"documents": processed_docs}
        )
        self.message_bus.publish(response)

class RetrievalAgent:
    """Agent responsible for embedding and semantic retrieval"""
    
    def __init__(self, message_bus: MessageBus):
        self.name = "RetrievalAgent"
        self.message_bus = message_bus
        self.message_bus.subscribe(self.name, self.handle_message)
        self.vector_store = None
    
    def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        if message.type == "INGESTION_COMPLETE":
            self.create_vector_store(message)
        elif message.type == "RETRIEVAL_REQUEST":
            self.retrieve_context(message)
    
    def create_vector_store(self, message: MCPMessage):
        """Create vector store from processed documents"""
        documents = message.payload.get("documents", [])
        
        if documents:
            try:
                self.vector_store = FAISS.from_documents(documents, embeddings)
                logger.info(f"Vector store created with {len(documents)} documents")
                
                # Notify completion
                response = MCPMessage(
                    sender=self.name,
                    receiver="CoordinatorAgent",
                    msg_type="VECTORSTORE_READY",
                    trace_id=message.trace_id,
                    payload={"status": "ready"}
                )
                self.message_bus.publish(response)
            except Exception as e:
                logger.error(f"Error creating vector store: {e}")
    
    def retrieve_context(self, message: MCPMessage):
        """Retrieve relevant context for a query"""
        query = message.payload.get("query", "")
        k = message.payload.get("k", 3)
        
        if self.vector_store and query:
            try:
                docs = self.vector_store.similarity_search(query, k=k)
                context = [{"content": doc.page_content, "source": doc.metadata.get("source", "")} 
                           for doc in docs]
                
                response = MCPMessage(
                    sender=self.name,
                    receiver="LLMResponseAgent",
                    msg_type="CONTEXT_RESPONSE",
                    trace_id=message.trace_id,
                    payload={
                        "query": query,
                        "retrieved_context": context,
                        "top_chunks": [doc.page_content for doc in docs]
                    }
                )
                self.message_bus.publish(response)
            except Exception as e:
                logger.error(f"Error retrieving context: {e}")

class LLMResponseAgent:
    """Agent responsible for generating LLM responses"""
    
    def __init__(self, message_bus: MessageBus):
        self.name = "LLMResponseAgent"
        self.message_bus = message_bus
        self.message_bus.subscribe(self.name, self.handle_message)
    
    def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        if message.type == "CONTEXT_RESPONSE":
            self.generate_response(message)
    
    def generate_response(self, message: MCPMessage):
        """Generate response using retrieved context"""
        query = message.payload.get("query", "")
        context = message.payload.get("retrieved_context", [])
        
        # Build context string
        context_text = "\n\n".join([f"Source: {ctx['source']}\nContent: {ctx['content']}" 
                                      for ctx in context])
        
        # Create messages for conversational format
        messages = [
            {
                "role": "system",
                "content": "You are a helpful assistant. Based on the provided context below, answer the user's question accurately and comprehensively. Cite the sources if possible.",
            },
            {
                "role": "user", 
                "content": f"Context:\n\n{context_text}\n\nQuestion: {query}"
            }
        ]

        try:
            # Use client.chat_completion for conversational models
            response_stream = client.chat_completion(
                messages=messages,
                max_tokens=512,
                temperature=0.7,
                stream=True
            )
            
            # Send streaming response
            response = MCPMessage(
                sender=self.name,
                receiver="CoordinatorAgent",
                msg_type="LLM_RESPONSE_STREAM",
                trace_id=message.trace_id,
                payload={
                    "query": query,
                    "response_stream": response_stream,
                    "context": context
                }
            )
            self.message_bus.publish(response)
            
        except Exception as e:
            logger.error(f"Error generating response: {e}")
            # Send an error stream back
            error_msg = f"Error from LLM: {e}"
            def error_generator():
                yield error_msg
            
            response = MCPMessage(
                sender=self.name,
                receiver="CoordinatorAgent",
                msg_type="LLM_RESPONSE_STREAM",
                trace_id=message.trace_id,
                payload={"response_stream": error_generator()}
            )
            self.message_bus.publish(response)

class CoordinatorAgent:
    """Coordinator agent that orchestrates the entire workflow"""
    
    def __init__(self, message_bus: MessageBus):
        self.name = "CoordinatorAgent"
        self.message_bus = message_bus
        self.message_bus.subscribe(self.name, self.handle_message)
        self.current_response_stream = None
        self.vector_store_ready = False
    
    def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        if message.type == "VECTORSTORE_READY":
            self.vector_store_ready = True
        elif message.type == "LLM_RESPONSE_STREAM":
            self.current_response_stream = message.payload.get("response_stream")
    
    def process_files(self, files):
        """Process uploaded files"""
        if not files:
            return "No files uploaded."
        
        file_paths = [file.name for file in files]
        
        # Send ingestion request
        message = MCPMessage(
            sender=self.name,
            receiver="IngestionAgent",
            msg_type="INGESTION_REQUEST",
            payload={"files": file_paths}
        )
        self.message_bus.publish(message)
        
        return f"Processing {len(files)} files: {', '.join([os.path.basename(fp) for fp in file_paths])}"
    
    def handle_query(self, query: str, history: List) -> Generator[str, None, None]:
        """Handle user query and return streaming response"""
        if not self.vector_store_ready:
            yield "Please upload and process documents first."
            return
        
        # Send retrieval request
        message = MCPMessage(
            sender=self.name,
            receiver="RetrievalAgent",
            msg_type="RETRIEVAL_REQUEST",
            payload={"query": query}
        )
        self.message_bus.publish(message)
        
        # Wait for response and stream
        import time
        timeout = 20  # seconds
        start_time = time.time()
        
        while not self.current_response_stream and (time.time() - start_time) < timeout:
            time.sleep(0.1)
        
        if self.current_response_stream:
            try:
                # Stream tokens directly
                for chunk in self.current_response_stream:
                    # The token is in chunk.choices[0].delta.content for chat_completion
                    if hasattr(chunk, 'choices') and chunk.choices:
                        token = chunk.choices[0].delta.content
                        if token:
                            yield token
                    else:
                        # Fallback for different response format
                        if hasattr(chunk, 'token'):
                            yield chunk.token
                        elif isinstance(chunk, str):
                            yield chunk
            except Exception as e:
                yield f"Error streaming response: {e}"
            finally:
                self.current_response_stream = None # Reset for next query
        else:
            yield "Timeout: No response received from LLM agent."

# Initialize agents
ingestion_agent = IngestionAgent(message_bus)
retrieval_agent = RetrievalAgent(message_bus)
llm_response_agent = LLMResponseAgent(message_bus)
coordinator_agent = CoordinatorAgent(message_bus)

def create_interface():
    """Create ChatGPT-style Gradio interface"""
    
    with gr.Blocks(
        theme=gr.themes.Base(),
        css="""
        /* Dark theme styling */
        .gradio-container {
            background-color: #1a1a1a !important;
            color: #ffffff !important;
            height: 100vh !important;
            max-width: none !important;
            padding: 0 !important;
        }
        
        /* Main container */
        .main-container {
            display: flex;
            flex-direction: column;
            height: 100vh;
            background: linear-gradient(135deg, #1a1a1a 0%, #2d2d2d 100%);
        }
        
        /* Header */
        .header {
            background: rgba(255, 193, 7, 0.1);
            border-bottom: 1px solid rgba(255, 193, 7, 0.2);
            padding: 1rem 2rem;
            backdrop-filter: blur(10px);
        }
        
        .header h1 {
            color: #ffc107;
            margin: 0;
            font-size: 1.5rem;
            font-weight: 600;
        }
        
        .header p {
            color: #cccccc;
            margin: 0.25rem 0 0 0;
            font-size: 0.9rem;
        }
        
        /* Chat area - REDUCED HEIGHT */
        .chat-container {
            flex: 1;
            display: flex;
            flex-direction: column;
            max-width: 1000px;
            margin: 0 auto;
            width: 100%;
            padding: 1rem;
            height: calc(100vh - 200px) !important; /* Reduced height */
        }
        
        /* Chatbot styling - SMALLER */
        .gradio-chatbot {
            height: 300px !important; /* Reduced from 500px */
            max-height: 300px !important;
            background: transparent !important;
            border: none !important;
            margin-bottom: 1rem;
            overflow-y: auto !important;
            box-shadow: 0 0 12px rgba(255, 193, 7, 0.1);
            
        }
        
        /* Input area */
        .input-area {
            background: rgba(45, 45, 45, 0.6);
            border-radius: 16px;
            padding: 1rem;
            border: 1px solid rgba(255, 193, 7, 0.2);
            backdrop-filter: blur(10px);
            position: sticky;
            bottom: 0;
        }
        
        /* File upload */
        .upload-area {
            background: rgba(255, 193, 7, 0.05);
            border: 2px dashed rgba(255, 193, 7, 0.3);
            border-radius: 12px;
            padding: 1rem;
            margin-bottom: 1rem;
            transition: all 0.3s ease;
        }
        
        /* Buttons - YELLOW SEND BUTTON */
        .send-btn {
            background: linear-gradient(135deg, #ffc107 0%, #ff8f00 100%) !important;
            color: #000000 !important;
            border: none !important;
            border-radius: 8px !important;
            font-weight: 600 !important;
            min-height: 40px !important;
        }
        
        .primary-btn {
            background: linear-gradient(135deg, #ffc107 0%, #ff8f00 100%) !important;
            color: #000000 !important;
            border: none !important;
            border-radius: 8px !important;
            font-weight: 600 !important;
        }
        
        /* Text inputs */
        .gradio-textbox input, .gradio-textbox textarea {
            background: rgba(45, 45, 45, 0.8) !important;
            color: #ffffff !important;
            border: 1px solid rgba(255, 193, 7, 0.2) !important;
            border-radius: 8px !important;
        }
        
        /* Processing indicator */
        .processing-indicator {
            background: rgba(255, 193, 7, 0.1);
            border: 1px solid rgba(255, 193, 7, 0.3);
            border-radius: 8px;
            padding: 0.75rem;
            margin: 0.5rem 0;
            color: #ffc107;
            text-align: center;
        }
        
        /* Input row styling */
        .input-row {
            display: flex !important;
            gap: 10px !important;
            align-items: end !important;
        }
        
        /* Message input */
        .message-input {
            flex: 1 !important;
            min-height: 40px !important;
        }
        """,
        title="Agentic RAG Assistant"
    ) as iface:

        # Header
        with gr.Row():
            with gr.Column():
                gr.HTML("""
                <div class="header">
                    <h1>Agentic RAG Assistant</h1>
                    <p>Upload documents and ask questions - powered by Multi-Agent Architecture</p>
                </div>
                """)

        # Main layout with sidebar and chat
        with gr.Row():
            # Left sidebar for file upload
            with gr.Column(scale=1):
                gr.Markdown("### 📁 Document Upload")
                
                file_upload = gr.File(
                    file_count="multiple",
                    file_types=[".pdf", ".pptx", ".csv", ".docx", ".txt", ".md"],
                    label="Upload Documents",
                    elem_classes=["upload-area"]
                )
                
                processing_status = gr.HTML(visible=False)
                
                process_btn = gr.Button(
                    "Process Documents",
                    variant="primary",
                    elem_classes=["primary-btn"]
                )
                
                # gr.Markdown("### ℹ️ Architecture")
                # gr.Markdown("""
                # **Multi-Agent System:**
                # - 📄 **IngestionAgent**: Document parsing
                # - 🔍 **RetrievalAgent**: Semantic search  
                # - 🤖 **LLMAgent**: Response generation
                # - 🎯 **CoordinatorAgent**: Workflow orchestration
                
                # **Features:**
                # - Streaming responses
                # - Multi-format support
                # - Context-aware answers
                # """)

            # Right side - Chat interface
            with gr.Column(scale=2):
                gr.Markdown("### 💬 Chat Interface")
                
                # Chatbot with reduced height
                chatbot = gr.Chatbot(
                    height=300,  # Reduced height
                    elem_classes=["gradio-chatbot"],
                    show_copy_button=True,
                    type="messages",
                    placeholder="Upload documents first, then start chatting!"
                )
                
                # Input area with improved layout
                with gr.Row(elem_classes=["input-row"]):
                    msg_input = gr.Textbox(
                        placeholder="Ask about your documents...",
                        label="Message",
                        scale=4,
                        elem_classes=["message-input"],
                        show_label=False,
                        autofocus=True
                    )
                    send_btn = gr.Button(
                        "Send", 
                        scale=1, 
                        elem_classes=["send-btn"],
                        size="sm"
                    )

                # Examples
                gr.Examples(
                    examples=[
                        "What are the main topics discussed?",
                        "Summarize the key findings",
                        "What metrics are mentioned?",
                        "What are the recommendations?"
                    ],
                    inputs=msg_input,
                    label="Example Questions"
                )

        # State to track document processing
        doc_processed = gr.State(False)
        
        # Event handlers
        def handle_file_upload_and_process(files):
            if not files:
                return gr.update(visible=False), False
            
            # Show processing indicator
            processing_html = f"""
            <div class="processing-indicator">
                📄 Processing {len(files)} documents... Please wait.
            </div>
            """
            
            # Process files
            try:
                result = coordinator_agent.process_files(files)
                
                # Wait a moment for processing to complete
                import time
                time.sleep(3)
                
                success_html = """
                <div style="background: rgba(76, 175, 80, 0.1); border: 1px solid rgba(76, 175, 80, 0.3); 
                           border-radius: 8px; padding: 0.75rem; color: #4caf50; text-align: center;">
                    Documents processed successfully! You can now ask questions.
                </div>
                """
                return gr.update(value=success_html, visible=True), True
                
            except Exception as e:
                error_html = f"""
                <div style="background: rgba(244, 67, 54, 0.1); border: 1px solid rgba(244, 67, 54, 0.3); 
                           border-radius: 8px; padding: 0.75rem; color: #f44336; text-align: center;">
                    ❌ Error processing documents: {str(e)}
                </div>
                """
                return gr.update(value=error_html, visible=True), False

        def respond(message, history, doc_ready):
            if not doc_ready:
                # Show error message
                history.append({"role": "user", "content": message})
                history.append({"role": "assistant", "content": " Please upload and process documents first."})
                return history, ""
            
            if not message.strip():
                return history, message
            
            # Add user message
            history.append({"role": "user", "content": message})
            history.append({"role": "assistant", "content": ""})
            
            # Stream response
            try:
                for token in coordinator_agent.handle_query(message, history):
                    history[-1]["content"] += token
                    yield history, ""
            except Exception as e:
                history[-1]["content"] = f"❌ Error: {str(e)}"
                yield history, ""

        # Event bindings
        process_btn.click(
            handle_file_upload_and_process,
            inputs=[file_upload],
            outputs=[processing_status, doc_processed]
        )

        send_btn.click(
            respond,
            inputs=[msg_input, chatbot, doc_processed],
            outputs=[chatbot, msg_input],
            show_progress=True
        )

        msg_input.submit(
            respond,
            inputs=[msg_input, chatbot, doc_processed],
            outputs=[chatbot, msg_input],
            show_progress=True
        )

    return iface

# Launch the application
if __name__ == "__main__":
    demo = create_interface()
    demo.launch(
        share=True,
        server_name="0.0.0.0",
        server_port=7860
    )
